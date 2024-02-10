from pptx.text import text

class _SettableRun(text.Subshape):
    """Text run object. Corresponds to ``<a:r>`` child element in a paragraph."""

    def __init__(self, r, parent):
        super(_SettableRun, self).__init__(parent)
        self._r = r

    @property
    def font(self):
        """
        |Font| instance containing run-level character properties for the
        text in this run. Character properties can be and perhaps most often
        are inherited from parent objects such as the paragraph and slide
        layout the run is contained in. Only those specifically overridden at
        the run level are contained in the font object.
        """
        rPr = self._r.get_or_add_rPr()
        return text.Font(rPr)

    @font.setter
    def font(self, font):
        self.font.bold = font.bold
        self.font.italic = font.italic
        self.font.language_id = font.language_id
        self.font.name = font.name
        self.font.size = font.size
        self.font.underline = font.underline
        if hasattr(font.color, 'rgb'):
            self.font.color.rgb = font.color.rgb

    @text.lazyproperty
    def hyperlink(self):
        """
        |_Hyperlink| instance acting as proxy for any ``<a:hlinkClick>``
        element under the run properties element. Created on demand, the
        hyperlink object is available whether an ``<a:hlinkClick>`` element
        is present or not, and creates or deletes that element as appropriate
        in response to actions on its methods and attributes.
        """
        rPr = self._r.get_or_add_rPr()
        return text._Hyperlink(rPr, self)

    @property
    def text(self):
        """Read/write. A unicode string containing the text in this run.

        Assignment replaces all text in the run. The assigned value can be a 7-bit ASCII
        string, a UTF-8 encoded 8-bit string, or unicode. String values are converted to
        unicode assuming UTF-8 encoding.

        Any other control characters in the assigned string other than tab or newline
        are escaped as a hex representation. For example, ESC (ASCII 27) is escaped as
        "_x001B_". Contrast the behavior of `TextFrame.text` and `_Paragraph.text` with
        respect to line-feed and vertical-tab characters.
        """
        return self._r.text

    @text.setter
    def text(self, str):
        self._r.text = text.to_unicode(str)


if __name__ == "__main__":
    from pptx import Presentation

    pptx_filename = 'example/nametag-example.pptx' 
    prs = Presentation(pptx_filename)

    slide = prs.slides[0]

    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_shapes.append(shape)
    assert len(text_shapes) >= 2
    font1 = text_shapes[0].text_frame.paragraphs[0].runs[0].font
    print(f"text1 size: {font1.size.pt}pt")
    
    font2 = text_shapes[1].text_frame.paragraphs[0].runs[0].font
    print(f"text2 size: {font2.size.pt}pt")
    
    font1 = font2
    print("after assign text2 to text1")
    print(f"text1 size: {font1.size.pt}pt")
   