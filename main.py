import io
import os
import sys
import json
import logging
import argparse

from dataclasses import dataclass
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tkinter import filedialog, messagebox

from src.draw_slide import SlideDrawer
from src.utils import get_data_by_sample, open_file_with_default_program, read_excel_data
from src.gui import get_args_by_gui

logging.basicConfig(level=logging.WARNING, format='%(levelname)s: %(message)s')

# TODO: Sample num(0부터 시작) -> Slide num(1부터 시작)으로 재정의
# TODO: pad(tag 옆 공간), margin(tag 간 공간) -> pad(tag 간 공간), margin(page 윗 공간) 재정의 
# TODO: label template 생성; 직사각현 실선으로 내용물 범위 표시, picture placeholder로 배경범위 표시
# TODO: label template 선택 UI 추가 +alpha 사용자 지정 가능하도록

def get_args():
    parser = argparse.ArgumentParser(description="Create nametag pptx from excel file")
    parser.add_argument("--excel", type=str, help="Excel file name")
    parser.add_argument("--pptx", type=str, help="PowerPoint file name")
    parser.add_argument("--margin_x", type=float, default=0.0, help="Margin between nametags in x direction. unit: cm")
    parser.add_argument("--margin_y", type=float, default=0.0, help="Margin between nametags in y direction. unit: cm")
    parser.add_argument("--padding_x", type=float, default=0.0, help="Padding of nametag in x direction. unit: cm")
    parser.add_argument("--padding_y", type=float, default=0.0, help="Padding of nametag in y direction. unit: cm")
    parser.add_argument("--per_slide", type=int, help="Number of nametags per slide")
    parser.add_argument("--gui", action="store_true", help="Use Tkinter GUI to select files and set parameters")
    parser.add_argument("--rpc", action="store_true", help="Pass arguments through JSON (Electron IPC)") # RPC: Remote Procedure Call
    args = parser.parse_args()

    if args.rpc:
        args.gui = True
    elif not args.excel or not args.pptx or args.gui:
        args = get_args_by_gui(args)
        args.gui = True
    return args

def assert_file_valid(file_path):
    if not file_path:
        raise ValueError("Missing required parameters")
    
    if not os.path.exists(file_path):
        raise ValueError(
            f"File Path Encoding Error\n"
            f"Problematic File: {file_path}\n"
            f"System Encodings:\n"
            f"  - Default: {sys.getdefaultencoding()}\n"
            f"  - Filesystem: {sys.getfilesystemencoding()}\n"
        )

@dataclass
class GetExcelHeaderRequest:
    excel: str
    
    def __post_init__(self):
        assert_file_valid(self.excel)

@dataclass
class GetPptxTextRequest:
    pptx: str
    
    def __post_init__(self):
        assert_file_valid(self.pptx)

@dataclass
class GenerateRequest:
    pptx: str
    excel: str
    margin_x: float = 0.0
    margin_y: float = 0.0
    padding_x: float = 0.0
    padding_y: float = 0.0
    per_slide: int = None

    def __post_init__(self):
        assert_file_valid(self.pptx)
        assert_file_valid(self.excel)

class TaskManger:
    def __init__(self, is_gui):
        self.is_gui = is_gui
        self.tasks = {
            "get_excel_header": (self.get_excel_header, GetExcelHeaderRequest),
            "get_pptx_slide_text": (self.get_pptx_slide_text, GetPptxTextRequest),
            "generate_pptx": (self.generate_pptx, GenerateRequest),
        }

    def process_request(self, request):
        task = request.get("task")
        data = request.get("data", {})

        response = {"task": task}
        if task in self.tasks:
            function, dataclass_type = self.tasks[task]
            try:
                request_data = dataclass_type(**data)
            except ValueError as e:
                response.update({"status": "error", "message": str(e)})
            except TypeError as e:
                response.update({"status": "developer_error", "message": f"Invalid parameters for {task}: {str(e)}"})
            else:
                _response = function(request_data)
                response.update(_response)
        else:
            response.update({"status": "developer_error", "message": f"Unknown task: {task}"})

        print(json.dumps(response, ensure_ascii=False), flush=True)
    
    def get_excel_header(self, data: GetExcelHeaderRequest):
        headers, _ = read_excel_data(data.excel)
        return {"status": "success", "headers": headers}

    def get_pptx_slide_text(self, data: GetPptxTextRequest):
        prs = Presentation(data.pptx)
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            shapes_to_process = list(slide.shapes)
            for shape in shapes_to_process:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text.strip())
                elif hasattr(shape, 'shapes') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    shapes_to_process.extend(shape.shapes)
            slides_text.append(slide_text)
        return {"status": "success", "slides": slides_text}
    
    def generate_pptx(self, data: GenerateRequest):
        prs = Presentation(data.pptx)
        sample_num = len(prs.slides)

        data_by_sample = get_data_by_sample(data.excel)
        for i in data_by_sample.keys():
            if isinstance(i, str):
                self.log_warning(f"Sample number '{i}' is not an integer. Skipping sample '{i}'.")
                continue
            if i >= sample_num:
                self.log_warning(f"No sample slide with index {i} exists. Skipping sample {i}.\nHint: The 'Sample Num' column should start from 0 and be continuous.")
                continue
            try:
                SlideDrawer(prs, i, data_by_sample[i]).draw(
                    margin=(data.margin_x, data.margin_y),
                    padding=(data.padding_x, data.padding_y),
                    per_slide=data.per_slide
                )
            except Exception as e:
                return {"status": "developer_error", "message": f"Error while drawing slide {i}: {str(e)}"}

        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")],
            title="Save the file as",
            initialfile=f"generated-{os.path.basename(data.pptx)}"
        )
        if filename:
            try:
                prs.save(filename)
                open_file_with_default_program(filename)
                return {"status": "success", "message": f"PPTX saved as '{os.path.basename(filename)}'"}
            except PermissionError:
                return {"status": "error", "message": f"Close the file '{os.path.basename(filename)}' to save"}
        else:
            return {"status": "success", "message": "Saving PPTX canceled by user"}

    def log_warning(self, message):
        if self.is_gui:
            messagebox.showwarning("Warning", message)
        logging.warning(message)


if __name__ == "__main__":
    args = get_args()
    task_manager = TaskManger(args.gui)

    if args.rpc:
        print("Python RPC mode ready")

        sys.stdin = io.TextIOWrapper(sys.stdin.buffer, encoding='utf-8')
        sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)
        sys.stderr.reconfigure(encoding='utf-8')
        
        for line in sys.stdin:
            try:
                request = json.loads(line.strip())
                task_manager.process_request(request)
            except Exception as e:
                print(json.dumps({"status": "developer_error", "message": str(e)}, ensure_ascii=False), flush=True)
    else:
        data = GenerateRequest(
            pptx=args.pptx,
            excel=args.excel,
            margin_x=args.margin_x,
            margin_y=args.margin_y,
            padding_x=args.padding_x,
            padding_y=args.padding_y,
            per_slide=args.per_slide
        )
        result = task_manager.generate_pptx(data)
        print(result["message"])
