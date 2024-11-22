import json
import argparse

from pptx import Presentation
from tkinter import filedialog

from src.utils import open_file_with_default_program
from pptx.util import Cm

def get_args():
    parser = argparse.ArgumentParser(description='Create nametag pptx from excel file')
    parser.add_argument('--label_name', type=str, help='Label name')
    args = parser.parse_args()

    return args

if __name__ == "__main__":
    args = get_args()

    ppt = "template/nametag.pptx"

    prs = Presentation(ppt)

    sample_slide_layout = 1
    slide_layout = prs.slide_layouts[sample_slide_layout]
    slide = prs.slides.add_slide(slide_layout)
    
    with open("template/label_templates.json", "r") as file:
        label_templates = json.load(file)
    
        for label_name in label_templates:
            if args.label_name in label_name:
                label = label_templates[label_name]
                assert label["page_size"] == "A4", f"Error: Page size '{label['page_size']}' is not supported"
                break
        else:
            raise ValueError(f"Error: Label name '{args.label_name}' is not found in label_templates.json")
    
    placeholder = slide.shapes[0]
    placeholder.width = Cm(label["label_width"])
    placeholder.height = Cm(label["label_height"])
    placeholder.left = Cm(prs.slide_width.cm / 2 - placeholder.width.cm / 2)
    placeholder.top = Cm(prs.slide_height.cm / 2 - placeholder.height.cm / 2)
    del prs.slides._sldIdLst[0]

    filename = filedialog.asksaveasfilename(
        defaultextension=".pptx",
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Save the file as",
        initialfile=f"{label['company']} {label['name']}.pptx"
    )
    if filename:
        prs.save(filename)
        print(f"Done: '{filename}' is saved")
        open_file_with_default_program(filename)
    else:
        print("Canceled by user.")