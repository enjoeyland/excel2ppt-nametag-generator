from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture

from src.utils import dotdict

def get_shape_info(shape: Picture|BaseShape):
    if isinstance(shape, Picture):
        return get_image_info(shape)
    elif isinstance(shape, Shape):
        return get_textbox_info(shape) or get_auto_shape_info(shape)
    else: # isinstance(shape, BaseShape):
        raise ValueError(f"Unsupported shape type: {type(shape)}")

def get_auto_shape_info(shape: Shape):
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        shape.fill.solid()
        return dotdict({
            "type": "auto_shape", 
            "left": shape.left.cm,
            "top": shape.top.cm,
            "width" : shape.width.cm,
            "height": shape.height.cm,
            "text": shape.text,
            "fill": shape.fill.fore_color,
            "line": shape.line.color,
            "shadow": shape.shadow,
            "shape_type": shape.shape_type
        })

def get_image_info(shape: Picture):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
        return dotdict({
            "type": "image", 
            "left": shape.left.cm,
            "top": shape.top.cm,
            "width" : shape.width.cm,
            "height": shape.height.cm,
            "image": shape.image,
            "crop": (shape.crop_left, shape.crop_right, shape.crop_top, shape.crop_bottom)
        })

def get_textbox_info(shape: Shape):
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return dotdict({
            "type": "text", 
            "left": shape.left.cm,
            "top": shape.top.cm,
            "width" : shape.width.cm,
            "height": shape.height.cm,
            "text": shape.text_frame.text,
            "alignment": shape.text_frame.paragraphs[0].alignment or PP_ALIGN.LEFT,
            "font": shape.text_frame.paragraphs[0].runs[0].font
        })

def to_relative_info(absolute_info):
    result = dotdict({
        "type": "layout",
        "images": [],
        "texts": [],
        "auto_shapes": [],
    })
    basis =  absolute_info[0].copy()
    for info in absolute_info:
        if info.left < basis.left and info.top < basis.top:
            basis =  info.copy()

    right = basis.left + basis.width
    buttom = basis.top + basis.height
    for info in absolute_info:
        right = max(right, info.left + info.width)
        buttom = max(buttom, info.top + info.height)

        info.left -= basis.left
        info.top -= basis.top
        if info.type == "text":
            result.texts.append(info)
        elif info.type == "image":
            result.images.append(info)
        elif info.type == "auto_shape":
            result.auto_shapes.append(info)
    
    result.width = right - basis.left
    result.height = buttom - basis.top
    return result

def get_slide_info(slide: Slide):
    absolute_info = []
    for shape in slide.shapes:
        print("shape:", shape.name, shape.shape_type)
        absolute_info.append(get_shape_info(shape))
    return to_relative_info(absolute_info)

if __name__ == "__main__":
    from pptx import Presentation

    pptx_filename = 'example/nametag-example.pptx' 
    prs = Presentation(pptx_filename)
    
    for i, slide in enumerate(prs.slides):
        print(f"{'='*5}#{i} slide{'='*5}")
        print(get_slide_info(slide))
