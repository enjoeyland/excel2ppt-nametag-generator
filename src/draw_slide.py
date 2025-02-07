from pptx.presentation import Presentation

from .draw_nametag import NameTagDrawer
from .utils import chunk_list

class SlidePositioner:
    def __init__(self, slide_size, sample, data, padding = (0, 0), margin = (0, 0), per_slide = None):
        self._slide_width, self._slide_height = slide_size
        self._sample = sample

        self.padding = padding
        self.margin = margin

        self.num_col, self.num_row = self.get_max_col_row()
        self.num_per_slide = self.num_col * self.num_row if per_slide is None else min(per_slide, self.num_col * self.num_row)
        self.left, self.top = self._get_start_pos()
        self.data_by_slide = chunk_list(data, self.num_per_slide)


    def get_max_col_row(self):
        num_col = (self._slide_width + self.margin[0]) / (self._sample.width + self.padding[0] * 2 + self.margin[0])
        num_row = (self._slide_height + self.margin[1]) / (self._sample.height + self.padding[1] * 2 + self.margin[1])
        return (int(num_col), int(num_row))

    def _get_start_pos(self):
        left = (self._slide_width - self.num_col * (self._sample.width + self.padding[0] * 2 + self.margin[0]) + self.margin[0]) / 2
        top = (self._slide_height - self.num_row * (self._sample.height + self.padding[1] * 2 + self.margin[1]) + self.margin[1]) / 2
        return (left, top)
    
    def _get_index(self, idx):
        col_idx = idx % self.num_col
        row_idx = idx // self.num_col
        return (col_idx, row_idx)
    
    def _get_position(self, idx):
        col_idx, row_idx = self._get_index(idx)
        return (
            self.left + col_idx * (self._sample.width + self.padding[0] * 2 + self.margin[0]) + self.padding[0],
            self.top + row_idx * (self._sample.height + self.padding[1] * 2 + self.margin[1]) + self.padding[1]
        )

    def slide_info_generator(self):
        for data in self.data_by_slide:
            yield self._nametag_info_generator(data)

    def _nametag_info_generator(self, slide_info):
        for idx, d in enumerate(slide_info):
            yield *self._get_position(idx), d

class SlideDrawer:
    def __init__(self, prs: Presentation, sample_num: int, data: list[dict[str, int|str]], blank_slide_layout = 0):
        self._prs = prs  
        self.data = data
        self.sample = NameTagDrawer.create_from_slide(self._prs.slides[sample_num])
        # self.slide_layout = self._prs.slides[sample_num].slide_layout
        self.slide_layout = self._prs.slide_layouts[blank_slide_layout]

    def draw(self, **kwargs):
        self.position = SlidePositioner((self._prs.slide_width.cm, self._prs.slide_height.cm), self.sample, self.data, **kwargs)

        for slide_info in self.position.slide_info_generator():
            slide = self._prs.slides.add_slide(self.slide_layout)
            for left, top, data in slide_info:
                self.sample.draw(slide, left, top)
                self.sample.set_text(data)
        return self._prs