from io import BytesIO
from abc import ABC, abstractmethod

from pptx.slide import Slide
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.shapes.connector import Connector
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Cm

from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.oxml.ns import qn

from .utils import set_fill, set_line, set_base_shape

class ShapeDrawer(ABC):
    def __init__(self, shape: Picture|BaseShape):
        self.shape = shape
        self.to_relative_position(0, 0)
        self.drawed_shape = None

    def to_relative_position(self, left: float, top: float):
        self.left:float = self.shape.left.cm - left
        self.top:float = self.shape.top.cm - top
    
    @abstractmethod
    def draw(self, slide: Slide, left: float=0, top: float=0):
        pass
    
    @classmethod
    def create(cls, shape: BaseShape):
        if isinstance(shape, GroupShape):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                return shape.shapes
        elif isinstance(shape, Picture):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                return ImageDrawer(shape)
        elif isinstance(shape, Connector):
            return ConnectorDrawer(shape)
        elif isinstance(shape, Shape):
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return TextBoxDrawer(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return AutoShapeDrawer(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                return None
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return None
        raise ValueError(f"Unsupported shape type: {type(shape)}")

class ImageDrawer(ShapeDrawer):
    def __init__(self, shape: Picture):
        super().__init__(shape)

    def draw(self, slide: Slide, left: float=0, top: float=0):
        shapes: SlideShapes = slide.shapes

        pic = shapes.add_picture(
            BytesIO(self.shape.image.blob),
            Cm(left + self.left),
            Cm(top + self.top),
            self.shape.width,
            self.shape.height
        )
        set_base_shape(self.shape, pic)
        set_line(self.shape.line, pic.line)
        pic.crop_left = self.shape.crop_left
        pic.crop_right = self.shape.crop_right
        pic.crop_top = self.shape.crop_top
        pic.crop_bottom = self.shape.crop_bottom
        self.drawed_shape = pic
        return pic

class TextBoxDrawer(ShapeDrawer):
    def __init__(self, shape: Shape):
        self.label = shape.text.strip().lower()
        super().__init__(shape)

    def draw(self, slide: Slide, left: float=0, top: float=0):
        shapes: SlideShapes = slide.shapes

        shape = shapes.add_textbox(
            Cm(left + self.left),
            Cm(top + self.top),
            self.shape.width,
            self.shape.height
        )
        set_base_shape(self.shape, shape)
        set_fill(self.shape, shape)
        set_line(self.shape.line, shape.line)
        p = shape.text_frame.paragraphs[0]
        p.alignment = self.shape.text_frame.paragraphs[0].alignment or PP_ALIGN.LEFT
        r = p.add_run()
        r.font = self.shape.text_frame.paragraphs[0].runs[0].font
        r.text = self.shape.text
        self.drawed_shape = p
        return p
    
    def set_text(self, text: str):
        if self.drawed_shape is None:
            raise ValueError("Shape is not drawn yet")
        self.drawed_shape.runs[0].text = text

class AutoShapeDrawer(ShapeDrawer):
    def __init__(self, shape: Shape):
        self.label = shape.text.strip().lower()
        super().__init__(shape)

    def draw(self, slide: Slide, left: float=0, top: float=0):
        shapes: SlideShapes = slide.shapes

        shape = shapes.add_shape(
            self.shape.auto_shape_type,
            Cm(left + self.left),
            Cm(top + self.top),
            self.shape.width,
            self.shape.height
        )
        set_base_shape(self.shape, shape)
        set_fill(self.shape, shape)
        set_line(self.shape.line, shape.line)
        shape.text = self.shape.text

        self.drawed_shape = shape
        return shape

    def set_text(self, text: str):
        if self.drawed_shape is None:
            raise ValueError("Shape is not drawn yet")
        self.drawed_shape.text = text

class ConnectorDrawer(ShapeDrawer):
    def __init__(self, shape: Picture):
        super().__init__(shape)

    def to_relative_position(self, left: float, top: float):
        self.begin_x = self.shape.begin_x.cm - left
        self.begin_y = self.shape.begin_y.cm - top
        self.end_x = self.shape.end_x.cm - left
        self.end_y = self.shape.end_y.cm - top
    
    def draw(self, slide: Slide, left: float=0, top: float=0):
        shapes: SlideShapes = slide.shapes
        connector = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Cm(left + self.begin_x),
            Cm(top + self.begin_y),
            Cm(left + self.end_x),
            Cm(top + self.end_y)
        )
        set_base_shape(self.shape, connector)
        set_line(self.shape.line, connector.line)
        self.set_arrow(connector.line, *self.get_arrow_type(self.shape.line))
        
        self.drawed_shape = connector
        return connector

    def get_arrow_type(self, line):
        line_elem = line._get_or_add_ln()

        start_arrow_elem = line_elem.find(f".//{qn('a:tailEnd')}")
        start_arrow = start_arrow_elem.get("type") if start_arrow_elem is not None else None
        
        end_arrow_elem = line_elem.find(f".//{qn('a:headEnd')}")
        end_arrow = end_arrow_elem.get("type") if end_arrow_elem is not None else None
        return start_arrow, end_arrow

    def set_arrow(self, line, start_arrow=None, end_arrow=None):
        line_elem = line._get_or_add_ln()

        for arrow in line_elem.findall(qn("a:headEnd")):
            line_elem.remove(arrow)
        for arrow in line_elem.findall(qn("a:tailEnd")):
            line_elem.remove(arrow)

        if start_arrow:
            line_elem.append(parse_xml(f'<a:tailEnd {nsdecls("a")} type="{start_arrow}"/>'))
        if end_arrow: # TODO: 반영이 안되고 있음...
            line_elem.append(parse_xml(f'<a:headEnd {nsdecls("a")} type="{end_arrow}"/>'))
