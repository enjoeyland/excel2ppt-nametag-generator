from io import BytesIO

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Cm

from src.analyze_slide import get_slide_info
from src.utils import chunk_list

class SlidePositioner:
    def __init__(self, slide_size, sample, data, padding = (0, 0), margin = (0, 0), per_slide = None):
        self._slide_width, self._slide_height = slide_size
        self._sample = sample

        self.padding = padding
        self.margin = margin

        self.num_col, self.num_row = self.get_max_col_row()
        self.num_per_slide = self.num_col * self.num_row if per_slide is None else min(per_slide, self.num_col * self.num_row)
        self.left, self.top = self._get_start_pos()
        self.data_by_slide = chunk_list(data, self.num_per_slide)


    def get_max_col_row(self):
        num_col = (self._slide_width + self.margin[0]) / (self._sample.width + self.padding[0] * 2 + self.margin[0])
        num_row = (self._slide_height + self.margin[1]) / (self._sample.height + self.padding[1] * 2 + self.margin[1])
        return (int(num_col), int(num_row))

    def _get_start_pos(self):
        left = (self._slide_width - self.num_col * (self._sample.width + self.padding[0] * 2 + self.margin[0]) + self.margin[0]) / 2
        top = (self._slide_height - self.num_row * (self._sample.height + self.padding[1] * 2 + self.margin[1]) + self.margin[1]) / 2
        return (left, top)
    
    def _get_index(self, idx):
        col_idx = idx % self.num_col
        row_idx = idx // self.num_col
        return (col_idx, row_idx)
    
    def _get_position(self, idx):
        col_idx, row_idx = self._get_index(idx)
        return (
            self.left + col_idx * (self._sample.width + self.padding[0] * 2 + self.margin[0]) + self.padding[0],
            self.top + row_idx * (self._sample.height + self.padding[1] * 2 + self.margin[1]) + self.padding[1]
        )

    def slide_info_generator(self):
        for data in self.data_by_slide:
            yield self._nametag_info_generator(data)

    def _nametag_info_generator(self, slide_info):
        for idx, d in enumerate(slide_info):
            yield *self._get_position(idx), d

class NameTagDrawer:
    def __init__(self, prs: Presentation, sample_num: int, data: list[dict[str, int|str]]):
        self._prs = prs  
        self.data = data
        self.sample = get_slide_info(self._prs.slides[sample_num])

    def add_nametag_info(self, slide: Slide, slide_info: list[tuple[int, int, dict[str, int|str]]]):
        shapes: SlideShapes  = slide.shapes
        for left, top, data in slide_info:
            for image_form in self.sample.images:
                pic = shapes.add_picture(
                    BytesIO(image_form.image.blob),
                    Cm(left + image_form.left),
                    Cm(top + image_form.top),
                    Cm(image_form.width),
                    Cm(image_form.height)
                )
                pic.crop_left, pic.crop_right, pic.crop_top, pic.crop_bottom = image_form.crop

            for text_form in self.sample.texts:
                p = shapes.add_textbox(
                    Cm(left + text_form.left),
                    Cm(top + text_form.top),
                    Cm(text_form.width),
                    Cm(text_form.height)
                ).text_frame.paragraphs[0]
                p.alignment = text_form.alignment
                r = p.add_run()
                r.font = text_form.font
                if text_form.text.lower() in data:
                    r.text = data[text_form.text.lower()]
                else:
                    r.text = text_form.text
            
            for auto_shape_form in self.sample.auto_shapes:
                shape = shapes.add_shape(
                    auto_shape_form.shape_type,
                    Cm(left + auto_shape_form.left),
                    Cm(top + auto_shape_form.top),
                    Cm(auto_shape_form.width),
                    Cm(auto_shape_form.height)
                )
                if auto_shape_form.text.lower() in data:
                    shape.text = data[auto_shape_form.text.lower()]
                else:
                    shape.text = auto_shape_form.text
                # TODO: change fill, line settable. see settable_pptx.py
                # shape.fill.solid()
                # print(type(auto_shape_form.fill))
                # shape.fill.fore_color.rgb = auto_shape_form.fill.rgb
                # shape.line.color.rgb = auto_shape_form.line.rgb

    def draw(self, blank_slide_layout = 0, **kwargs):
        self.position = SlidePositioner((self._prs.slide_width.cm, self._prs.slide_height.cm), self.sample, self.data, **kwargs)
        slide_layout = self._prs.slide_layouts[blank_slide_layout]

        for slide_info in self.position.slide_info_generator():
            slide = self._prs.slides.add_slide(slide_layout)
            self.add_nametag_info(slide, slide_info)
        return self._prs