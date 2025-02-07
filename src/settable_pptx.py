from pptx.text.text import _Run

from .utils import set_color

# Add setter for font
def font(self, font):
    self.font.bold = font.bold
    self.font.italic = font.italic
    self.font.language_id = font.language_id
    self.font.name = font.name
    self.font.size = font.size
    self.font.underline = font.underline
    set_color(font, self.font)

_Run.font = _Run.font.setter(font)

if __name__ == "__main__":
    from pptx import Presentation

    pptx_filename = 'example/nametag-example.pptx' 
    prs = Presentation(pptx_filename)

    slide = prs.slides[0]

    text_shapes = []
    for i, shape in  enumerate(slide.shapes):
        if shape.has_text_frame:
            text_shapes.append(i)
    assert len(text_shapes) >= 2


    font1 = slide.shapes[text_shapes[0]].text_frame.paragraphs[0].runs[0].font
    print(f"text1 size: {font1.size.pt}pt")
    
    font2 = slide.shapes[text_shapes[1]].text_frame.paragraphs[0].runs[0].font
    print(f"text2 size: {font2.size.pt}pt")

    slide.shapes[text_shapes[0]].text_frame.paragraphs[0].runs[0].font = font2
    print("after assign text2 to text1")
    print(f"text1 size: {font1.size.pt}pt")
   