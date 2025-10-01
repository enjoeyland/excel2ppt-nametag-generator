import pptx.presentation

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from dataclasses import dataclass
from typing import Literal


@dataclass
class LabelTemplate:
    company: str
    name: str
    category: str
    page_size: str
    horizontal_num: int
    vertical_num: int
    top_margin: float
    side_margin: float
    label_height: float
    label_width: float
    horizontal_gap: float
    vertical_gap: float

prs: pptx.presentation.Presentation = Presentation("template/nametag2.pptx")

SLIDE_SIZE = {
    "default": (Cm(25.4), Cm(19.05)),   # 기본 (16:9 비율)
    "A3": (Cm(35.56), Cm(26.67)),       # A3 (297mm x 420mm)
    "A4": (Cm(27.517), Cm(19.05)),      # A4 (210mm x 297mm)
    "Letter": (Cm(25.4), Cm(19.05)),    # Letter (8.5in x 11in)
    "Ledger": (Cm(33.831), Cm(25.374))  # Ledger (11in x 17in)
}

SlideSizeType = Literal["default", "A3", "A4", "Letter", "Ledger"]
SlideOrientation = Literal["landscape", "portrait"]

def set_slide_size(prs: Presentation, size: SlideSizeType = "default", orientation: SlideOrientation = "landscape"):
    if size not in SLIDE_SIZE:
        raise ValueError(f"'{size}'은(는) 지원되지 않는 슬라이드 크기입니다. 사용 가능한 옵션: {list(SLIDE_SIZE.keys())}")
    prs.slide_width, prs.slide_height = SLIDE_SIZE[size][:: -1] if orientation == "portrait" else SLIDE_SIZE[size]

def set_label_template(prs: Presentation, label_info: str) -> None:
    prs.core_properties.comments = label_info

def get_label_template(prs: Presentation) -> str:
    return prs.core_properties.comments or None


def create_template_slide(prs: Presentation, label_template: LabelTemplate) -> None:
    """
    'template/nametag.pptx' 파일에서 slide_layouts[1]을 사용하여 새로운 슬라이드를 추가하고,
    - 직사각형 실선으로 내용물 범위를 표시
    - Picture Placeholder를 찾아서 배경 영역으로 사용

    :param prs: pptx.presentation.Presentation 객체
    """
    # 📌 slide_layouts[1]을 사용 (이 레이아웃에는 Picture Placeholder가 포함됨)
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)

    # 🔹 기존 내용물 제거 (자리 표시자 제외)
    for shape in slide.shapes:
        if not shape.is_placeholder:
            slide.shapes._spTree.remove(shape._element)  # XML 요소 제거

    # 📏 직사각형 실선으로 내용물 범위 표시
    width, height = Cm(label_template.label_width), Cm(label_template.label_height)
    left, top = int((prs.slide_width - width)/2), int((prs.slide_height - height)/2)
    content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    content_box.fill.background() 
    content_box.line.color.rgb = RGBColor(0, 176, 240)  # 검정색 실선
    content_box.line.width = Pt(1)  # 선 두께: 1pt
    text_frame = content_box.text_frame
    text_frame.text = "재단선"
    text_frame.paragraphs[0].font.size = Pt(10)
    # text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY_LOW
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            placeholder = shape
            break

    placeholder.width, placeholder.height = int(Cm(label_template.label_width) + Cm(label_template.horizontal_gap)/2), int(Cm(label_template.label_height) + Cm(label_template.vertical_gap)/2)
    placeholder.left, placeholder.top = int((prs.slide_width - placeholder.width)/2), int((prs.slide_height - placeholder.height)/2)

def clear_slides(prs: Presentation) -> None:
    for i in range(len(prs.slides)-1, -1, -1):
        delete_slide(prs, i)

def delete_slide(prs: Presentation, idx: int) -> None:
    slide_count = len(prs.slides)
    if slide_count == 0:
        print("⚠️ 삭제할 슬라이드가 없습니다!")
        return
    if not (-slide_count <= idx < slide_count):
        print(f"⚠️ 유효하지 않은 인덱스입니다!")
        return
    
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

lt = LabelTemplate("FORMTEC", "3100", "바코드라벨", "A4", 5, 13, 1.07, 0.46, 2.12, 3.81, 0.25, 0.0)

clear_slides(prs)
set_slide_size(prs, lt.page_size, "portrait")
set_label_template(prs, f"{lt.company} {lt.name} {lt.category} {lt.page_size} {lt.horizontal_num}x{lt.vertical_num} margin: {lt.top_margin}, {lt.side_margin} label: {lt.label_height}, {lt.label_width} gap: {lt.horizontal_gap}, {lt.vertical_gap}")
create_template_slide(prs, lt)

prs.save(f"example/{lt.company}_{lt.name}_{lt.category}.pptx")
